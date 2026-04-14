"""PPTX parser companion for slide extraction."""

from __future__ import annotations

from typing import BinaryIO

from app.core.logging import LogManager

from .parser_contracts import DocumentParser, ParsedPage

logger = LogManager.get_logger("ai.rag.parser")


class PptxParser(DocumentParser):
    """
    PPTX Parser / PPTX 解析器

    Uses python-pptx to extract per-slide:
    - Text box/placeholder content
    - Tables (Markdown format)
    - Slide title (metadata.heading)
    - Slide notes (if any)
    使用 python-pptx 按幻灯片提取：
    - 文字框/占位符内容
    - 表格（Markdown 格式）
    - 幻灯片标题（metadata.heading）
    - 幻灯片备注（若有）

    Each slide → 1 ParsedPage (blank slides are skipped).
    每张幻灯片 → 1个 ParsedPage（空白幻灯片跳过）。
    """

    async def parse(
        self, file_content: BinaryIO, file_name: str = ""
    ) -> list[ParsedPage]:
        from pptx import Presentation

        prs = Presentation(file_content)
        pages: list[ParsedPage] = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            parts: list[str] = []
            title_text: str | None = None

            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if not text:
                        continue
                    if shape.shape_id == 1 or (
                        hasattr(shape, "placeholder_format")
                        and shape.placeholder_format is not None
                        and shape.placeholder_format.idx == 0
                    ):
                        title_text = text
                    parts.append(text)

                if shape.has_table:
                    rows_text = []
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        rows_text.append("| " + " | ".join(cells) + " |")
                    if rows_text:
                        header_sep = (
                            "| " + " | ".join(["---"] * len(shape.table.columns)) + " |"
                        )
                        table_md = (
                            rows_text[0]
                            + "\n"
                            + header_sep
                            + "\n"
                            + "\n".join(rows_text[1:])
                        )
                        parts.append(table_md)

            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    parts.append(notes_text)

            content = "\n\n".join(parts).strip()
            if not content:
                continue

            pages.append(
                ParsedPage(
                    content=content,
                    metadata={
                        "slide": slide_num,
                        "heading": title_text,
                        "source": file_name,
                    },
                )
            )

        logger.info(
            "PPTX parsed: {}, slides={}",
            file_name,
            len(pages),
        )
        return pages


__all__ = ["PptxParser"]
